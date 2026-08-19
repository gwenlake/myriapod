"""Write a swarm answer to a file: Markdown, PDF or PowerPoint.

Workers return Markdown, so the final answer already carries structure —
headings, bullets, emphasis, code. That structure is parsed once into a
small block model here, then rendered per format, so a PDF looks like a
document and a deck looks like a deck instead of one wall of text.

``.md`` and ``.txt`` need nothing beyond the standard library. ``.pdf``
and ``.pptx`` need the ``docs`` extra::

    uv add 'myriapod[docs]'
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable, Sequence

#: Extensions :func:`save_report` knows how to write.
FORMATS = (".md", ".markdown", ".txt", ".pdf", ".pptx")


class ReportError(RuntimeError):
    """Unsupported target format, or a missing optional dependency."""


# --------------------------------------------------------------------- #
# Markdown -> block model
# --------------------------------------------------------------------- #


@dataclass
class Span:
    """An inline run of text with its emphasis."""

    text: str
    bold: bool = False
    italic: bool = False
    code: bool = False


@dataclass
class Block:
    """One block-level element of the answer.

    ``kind`` is ``heading``, ``para``, ``bullet``, ``code``, ``rule`` or
    ``table``. ``level`` is the heading level (1-6) or the bullet nesting
    depth (0+). ``rows`` holds a table's cells, header row first.
    """

    kind: str
    spans: list[Span] = field(default_factory=list)
    level: int = 0
    text: str = ""  # raw text, code blocks only
    rows: list[list[list[Span]]] = field(default_factory=list)  # tables only

    @property
    def plain(self) -> str:
        if self.kind == "code":
            return self.text
        if self.kind == "table":
            return "\n".join(
                " | ".join("".join(s.text for s in cell) for cell in row)
                for row in self.rows
            )
        return "".join(s.text for s in self.spans)


_FENCE = re.compile(r"^\s*```")
_HEADING = re.compile(r"^(#{1,6})\s+(.*)$")
_RULE = re.compile(r"^\s*([-*_])(\s*\1){2,}\s*$")
_BULLET = re.compile(r"^(\s*)(?:[-*+]|\d+[.)])\s+(.*)$")
_QUOTE = re.compile(r"^\s*>\s?(.*)$")
# A table row is any line with a pipe inside it; the separator row under the
# header (|---|:--:|) is what confirms the block really is a table.
_TROW = re.compile(r"^\s*\|?(.*\|.*?)\|?\s*$")
_TSEP = re.compile(r"^\s*\|?[\s|]*:?-{2,}:?[\s|:-]*\|?\s*$")
_LINK = re.compile(r"\[([^\]]+)\]\((\S+?)\)")
# Only ** and * for emphasis: '_' is left alone so snake_case identifiers
# in worker output survive intact.
_INLINE = re.compile(r"(`[^`]+`|\*\*[^*]+\*\*|\*[^*]+\*)")


def parse_inline(text: str) -> list[Span]:
    """Split one line into emphasis runs. Links become ``text (url)``."""
    text = _LINK.sub(lambda m: f"{m.group(1)} ({m.group(2)})", text)
    spans: list[Span] = []
    for part in _INLINE.split(text):
        if not part:
            continue
        if part.startswith("`") and part.endswith("`") and len(part) > 1:
            spans.append(Span(part[1:-1], code=True))
        elif part.startswith("**") and part.endswith("**"):
            spans.append(Span(part[2:-2], bold=True))
        elif part.startswith("*") and part.endswith("*") and len(part) > 1:
            spans.append(Span(part[1:-1], italic=True))
        else:
            spans.append(Span(part))
    return spans or [Span("")]


def _table_cells(line: str) -> list[list[Span]]:
    """Split one table row into parsed cells, dropping the outer pipes."""
    return [parse_inline(cell.strip()) for cell in line.strip().strip("|").split("|")]


def parse_markdown(text: str) -> list[Block]:
    """Parse Markdown into blocks. Unknown syntax degrades to paragraphs."""
    blocks: list[Block] = []
    para: list[str] = []

    def flush() -> None:
        if para:
            blocks.append(Block("para", parse_inline(" ".join(para))))
            para.clear()

    lines = text.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i]

        if _FENCE.match(line):
            flush()
            i += 1
            body: list[str] = []
            while i < len(lines) and not _FENCE.match(lines[i]):
                body.append(lines[i])
                i += 1
            i += 1  # closing fence (or end of input)
            blocks.append(Block("code", text="\n".join(body)))
            continue

        # A pipe table only counts as one if the second line is its separator;
        # otherwise a sentence with a pipe in it would swallow what follows.
        if (
            "|" in line
            and i + 1 < len(lines)
            and _TSEP.match(lines[i + 1])
            and _TROW.match(line)
        ):
            flush()
            rows = [_table_cells(line)]
            i += 2  # header + separator
            while i < len(lines) and "|" in lines[i] and lines[i].strip():
                rows.append(_table_cells(lines[i]))
                i += 1
            width = max(len(r) for r in rows)
            for row in rows:
                row += [[] for _ in range(width - len(row))]
            blocks.append(Block("table", rows=rows))
            continue

        if not line.strip():
            flush()
        elif m := _HEADING.match(line):
            flush()
            blocks.append(Block("heading", parse_inline(m.group(2).strip()),
                                level=len(m.group(1))))
        elif _RULE.match(line):
            flush()
            blocks.append(Block("rule"))
        elif m := _BULLET.match(line):
            flush()
            blocks.append(Block("bullet", parse_inline(m.group(2).strip()),
                                level=len(m.group(1)) // 2))
        elif m := _QUOTE.match(line):
            flush()
            blocks.append(Block("para", parse_inline(m.group(1).strip())))
        else:
            para.append(line.strip())
        i += 1

    flush()
    return blocks


# --------------------------------------------------------------------- #
# Shared helpers
# --------------------------------------------------------------------- #

Meta = Sequence[tuple[str, str]]


def _meta_lines(meta: Meta | None) -> list[str]:
    return [f"{label}: {value}" for label, value in (meta or [])]


def _document_title(content: str, title: str | None) -> str | None:
    """The title to print, or None when the answer already carries its own.

    ``title`` is usually the goal as the user phrased it ("fais moi un topo
    sur..."), which reads badly stacked above the report's real title. A
    synthesis task almost always opens with one, so defer to it.
    """
    if not title:
        return None
    first = next((b for b in parse_markdown(content) if b.kind != "rule"), None)
    if first is not None and first.kind == "heading" and first.level == 1:
        return None
    return title


# --------------------------------------------------------------------- #
# Markdown / text
# --------------------------------------------------------------------- #


def write_markdown(
    path: Path, content: str, *, title: str | None = None, meta: Meta | None = None
) -> None:
    """Write the answer verbatim, with a provenance footer."""
    parts = []
    if title := _document_title(content, title):
        parts += [f"# {title}", ""]
    body = content.strip().splitlines()
    if meta:
        # The footer brings its own rule; a report that already signs off with
        # one would otherwise show two in a row.
        while body and _RULE.match(body[-1]):
            body.pop()
    parts.append("\n".join(body).rstrip())
    if meta:
        parts += ["", "---", "", "  \n".join(f"*{line}*" for line in _meta_lines(meta))]
    path.write_text("\n".join(parts) + "\n", encoding="utf-8")


def write_text(
    path: Path, content: str, *, title: str | None = None, meta: Meta | None = None
) -> None:
    """Plain text: the block model flattened, no Markdown punctuation."""
    out: list[str] = []
    if title := _document_title(content, title):
        out += [title, "=" * len(title), ""]
    for block in parse_markdown(content):
        if block.kind == "heading":
            out += ["", block.plain, "-" * len(block.plain)]
        elif block.kind == "bullet":
            out.append(f"{'  ' * block.level}- {block.plain}")
        elif block.kind == "rule":
            out.append("")
        elif block.kind == "code":
            out += [""] + [f"    {ln}" for ln in block.text.splitlines()] + [""]
        elif block.kind == "table":
            cells = [["".join(s.text for s in c) for c in row] for row in block.rows]
            widths = [max(len(r[i]) for r in cells) for i in range(len(cells[0]))]
            pad = lambda row: "  ".join(c.ljust(w) for c, w in zip(row, widths))
            out += ["", pad(cells[0]), "-" * len(pad(cells[0]))]
            out += [pad(row) for row in cells[1:]] + [""]
        else:
            out += ["", block.plain]
    if meta:
        out += ["", "---"] + _meta_lines(meta)
    path.write_text("\n".join(out).strip() + "\n", encoding="utf-8")


# --------------------------------------------------------------------- #
# PDF
# --------------------------------------------------------------------- #

#: Font families to try, as (regular, bold, italic) filenames. The first
#: family whose regular face exists wins; missing bold/italic fall back to
#: the regular face. Without a hit, the PDF uses the built-in Helvetica,
#: which is Latin-1 only (see ``_LATIN1_FALLBACK``).
_FONT_CANDIDATES: tuple[tuple[str, str, str], ...] = (
    (
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
        "/System/Library/Fonts/Supplemental/Arial Italic.ttf",
    ),
    (
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Oblique.ttf",
    ),
    (
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Italic.ttf",
    ),
    (
        "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/arialbd.ttf",
        "C:/Windows/Fonts/ariali.ttf",
    ),
)

#: Typographic characters the Latin-1 core fonts cannot encode.
_LATIN1_FALLBACK = str.maketrans({
    "\u2018": "'", "\u2019": "'", "\u201c": '"', "\u201d": '"',
    "\u2013": "-", "\u2014": "--", "\u2026": "...", "\u2022": "-",
    "\u00a0": " ", "\u202f": " ", "\u2192": "->", "\u2264": "<=",
    "\u2265": ">=", "\u2248": "~", "\u2713": "v", "\u2714": "v",
    "\u2717": "x", "\u2718": "x", "\u20ac": "EUR",
})


def _install_font(pdf) -> tuple[str, bool]:
    """Register a Unicode family; return (family, is_unicode)."""
    for regular, bold, italic in _FONT_CANDIDATES:
        if not Path(regular).exists():
            continue
        pdf.add_font("body", "", regular)
        pdf.add_font("body", "B", bold if Path(bold).exists() else regular)
        pdf.add_font("body", "I", italic if Path(italic).exists() else regular)
        return "body", True
    return "helvetica", False


def write_pdf(
    path: Path, content: str, *, title: str | None = None, meta: Meta | None = None
) -> None:
    """Render the answer as a paginated A4 document."""
    try:
        from fpdf import FPDF
        from fpdf.enums import XPos, YPos
        from fpdf.fonts import FontFace
    except ImportError as e:  # pragma: no cover - exercised without the extra
        raise ReportError(
            "Writing PDF needs the 'docs' extra: uv add 'myriapod[docs]'"
        ) from e

    def block_cell(height: float, text: str, **kwargs) -> None:
        """multi_cell that returns the cursor to the left margin.

        Its default lands the cursor at the *right* edge of the cell, which
        leaves the next full-width call zero space to render into.
        """
        pdf.multi_cell(0, height, text, new_x=XPos.LMARGIN, new_y=YPos.NEXT, **kwargs)

    pdf = FPDF(format="A4")
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.set_margins(20, 18, 20)
    family, unicode_ok = _install_font(pdf)
    mono = "courier"
    pdf.add_page()

    def clean(text: str) -> str:
        return text if unicode_ok else text.translate(_LATIN1_FALLBACK)

    def spans(runs: Iterable[Span], size: float, bold: bool = False) -> None:
        for span in runs:
            if span.code:
                pdf.set_font(mono, "", size - 0.5)
            else:
                style = ("B" if span.bold or bold else "") + ("I" if span.italic else "")
                pdf.set_font(family, style, size)
            pdf.write(size * 0.55, clean(span.text))
        pdf.ln(size * 0.62)

    if title := _document_title(content, title):
        pdf.set_font(family, "B", 22)
        block_cell(10, clean(title))
        pdf.ln(4)

    _H_SIZES = {1: 19.0, 2: 15.0, 3: 12.5, 4: 11.0, 5: 10.5, 6: 10.5}
    body = 10.5

    for block in parse_markdown(content):
        if block.kind == "heading":
            pdf.ln(3)
            spans(block.spans, _H_SIZES.get(block.level, body), bold=True)
            pdf.ln(1)
        elif block.kind == "bullet":
            indent = 20 + 5 * (block.level + 1)
            pdf.set_left_margin(indent)
            pdf.set_x(indent - 4)
            pdf.set_font(family, "", body)
            pdf.write(body * 0.55, clean("\u2022 "))
            spans(block.spans, body)
            pdf.set_left_margin(20)
        elif block.kind == "code":
            pdf.ln(1)
            pdf.set_font(mono, "", body - 1)
            pdf.set_fill_color(244, 245, 247)
            block_cell(body * 0.52, clean(block.text), fill=True)
            pdf.ln(2)
        elif block.kind == "table":
            # fpdf2 draws the grid, wraps the cells and repeats the header
            # across a page break; all it wants is plain strings, so inline
            # emphasis is flattened here rather than lost silently.
            pdf.ln(2)
            pdf.set_font(family, "", body - 1.5)
            head = FontFace(emphasis="BOLD", fill_color=(238, 240, 244))
            with pdf.table(headings_style=head, line_height=body * 0.5,
                           text_align="LEFT", padding=1.6) as table:
                for row in block.rows:
                    line = table.row()
                    for cell in row:
                        line.cell(clean("".join(s.text for s in cell)))
            pdf.ln(3)
        elif block.kind == "rule":
            pdf.ln(2)
            y = pdf.get_y()
            pdf.set_draw_color(200, 203, 208)
            pdf.line(20, y, pdf.w - 20, y)
            pdf.ln(3)
        else:
            spans(block.spans, body)
            pdf.ln(1.5)

    if meta:
        pdf.ln(4)
        pdf.set_draw_color(200, 203, 208)
        pdf.line(20, pdf.get_y(), pdf.w - 20, pdf.get_y())
        pdf.ln(3)
        pdf.set_font(family, "I", 8.5)
        pdf.set_text_color(110, 116, 126)
        for line in _meta_lines(meta):
            block_cell(4.5, clean(line))

    pdf.output(str(path))


# --------------------------------------------------------------------- #
# PowerPoint
# --------------------------------------------------------------------- #

#: Body items per slide before spilling onto a continuation slide.
_SLIDE_ITEMS = 8


def _flatten_tables(blocks: list[Block]) -> list[Block]:
    """One bullet per table row: a slide cannot host a grid and stay readable.

    The header row labels the fields, so it rides along inside each bullet
    ("Idée: audit-éclair · Prix: 4 900 €") rather than becoming a bullet of
    its own that reads as a heading with no content.
    """
    out: list[Block] = []
    for block in blocks:
        if block.kind != "table":
            out.append(block)
            continue
        head = ["".join(s.text for s in cell) for cell in block.rows[0]]
        for row in block.rows[1:]:
            cells = ["".join(s.text for s in cell) for cell in row]
            line = " · ".join(
                f"{label}: {value}" if label else value
                for label, value in zip(head, cells)
                if value
            )
            out.append(Block("bullet", parse_inline(line)))
    return out


def _slides(content: str, fallback_title: str) -> list[tuple[str, list[Block]]]:
    """Group blocks into (slide title, body blocks): h1/h2 start a slide."""
    out: list[tuple[str, list[Block]]] = []
    title, body = fallback_title, []
    for block in _flatten_tables(parse_markdown(content)):
        if block.kind == "rule":
            continue
        if block.kind == "heading" and block.level <= 2:
            if body or out:
                out.append((title, body))
            title, body = block.plain, []
        else:
            body.append(block)
    out.append((title, body))
    return [(t, b) for t, b in out if b or len(out) == 1]


def write_pptx(
    path: Path, content: str, *, title: str | None = None, meta: Meta | None = None
) -> None:
    """Render the answer as a 16:9 deck, one slide per top-level heading."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as e:  # pragma: no cover - exercised without the extra
        raise ReportError(
            "Writing PPTX needs the 'docs' extra: uv add 'myriapod[docs]'"
        ) from e

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    deck_title = title or "myriapod"

    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = deck_title
    if len(cover.placeholders) > 1:
        cover.placeholders[1].text = "\n".join(_meta_lines(meta)) or ""

    def add_paragraph(frame, block: Block, first: bool) -> None:
        para = frame.paragraphs[0] if first else frame.add_paragraph()
        if block.kind == "heading":
            para.level = 0
        elif block.kind == "bullet":
            para.level = min(block.level + 1, 4)
        else:
            para.level = 0
        runs = (
            [Span(block.text, code=True)] if block.kind == "code" else block.spans
        )
        for span in runs:
            run = para.add_run()
            run.text = span.text
            run.font.bold = span.bold or block.kind == "heading"
            run.font.italic = span.italic
            run.font.size = Pt(18 if block.kind == "heading" else 16)
            if span.code:
                run.font.name = "Menlo"

    for slide_title, blocks in _slides(content, deck_title):
        for start in range(0, max(len(blocks), 1), _SLIDE_ITEMS):
            chunk = blocks[start:start + _SLIDE_ITEMS]
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = (
                slide_title if start == 0 else f"{slide_title} (suite)"
            )
            frame = slide.placeholders[1].text_frame
            frame.word_wrap = True
            frame.clear()
            for n, block in enumerate(chunk):
                add_paragraph(frame, block, first=(n == 0))

    prs.save(str(path))


# --------------------------------------------------------------------- #
# Dispatch
# --------------------------------------------------------------------- #

_WRITERS = {
    ".md": write_markdown,
    ".markdown": write_markdown,
    ".txt": write_text,
    ".pdf": write_pdf,
    ".pptx": write_pptx,
}


def save_report(
    path: str | Path,
    content: str,
    *,
    title: str | None = None,
    meta: Meta | None = None,
) -> Path:
    """Write ``content`` to ``path``, choosing the format from its suffix.

    Args:
        path: Target file; the suffix picks the writer (see :data:`FORMATS`).
        content: The swarm's final answer, as Markdown.
        title: Document title / first slide, usually the goal.
        meta: ``(label, value)`` pairs recorded as a provenance footer.

    Raises:
        ReportError: Unknown suffix, or the ``docs`` extra is missing.
    """
    path = Path(path)
    writer = _WRITERS.get(path.suffix.lower())
    if writer is None:
        raise ReportError(
            f"Cannot write {path.suffix or 'a file with no extension'!r}. "
            f"Supported: {', '.join(FORMATS)}."
        )
    path.parent.mkdir(parents=True, exist_ok=True)
    writer(path, content, title=title, meta=meta)
    return path
